r"""The supervisor-facing slide deck (.pptx).

Run this file directly.

    python scripts/make_presentation.py

Writes `Phase_0-4_Review.pptx` next to the figures. One slide per figure, each
answering the same four questions: what it shows, the setup behind it, the data
it was computed from, and why it matters.

**Layout adapts to the figure.** Aspect ratios here run from 0.91 (a tall column
of digit samples) to 3.36 (a three-panel heatmap), so a single template would
either crop the wide ones or strand the tall ones in whitespace. Wide figures get
the image on top and the notes underneath; tall and square ones get the image
left and the notes in a right-hand column.

**Numbers in the notes are typed, not computed** — unlike the .docx summary. The
slide text is interpretive ("this is why the star is an outlier"), and the
figures themselves carry the values. Where a number does appear it is one already
stated in `docs/results.md`.
"""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "src"))

FIGURES = Path(r"C:\Users\alter\OneDrive\Desktop\PhD\Distributed Online Learning\preliminary work")

# --- the same palette the figures use, so the deck does not clash with them ---
SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK_SECONDARY = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x89, 0x87, 0x81)
ACCENT = RGBColor(0x2A, 0x78, 0xD6)
RULE = RGBColor(0xE1, 0xE0, 0xD9)

WIDE, TALL = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.55)


# --------------------------------------------------------------------------- #
# slide content
# --------------------------------------------------------------------------- #

Slide = dict


def sections() -> list[tuple[str, str, list[Slide]]]:
    """(section title, subtitle, slides)."""
    return [
        (
            "The benchmark",
            "What the agents see, and the knobs that define a run",
            [
                dict(
                    file="01_topologies.png",
                    title="Communication topologies",
                    claim="Eight graphs spanning three orders of magnitude in connectivity.",
                    shows="Each graph on N = 10 agents, with its spectral gap and diameter. "
                    "`disconnected` is a deliberate negative control.",
                    setup="Metropolis mixing weights, which are symmetric and doubly "
                    "stochastic — so the spectral gap is well defined for every one.",
                    data="Constructed, not measured. Built by env/graph.py at run start.",
                    why="Connectivity is the axis of research question Q1. Everything about "
                    "who can talk to whom is fixed here.",
                ),
                dict(
                    file="03_mnist_samples.png",
                    title="The task",
                    claim="Ten-class MNIST, the standard benchmark, deliberately unambitious.",
                    shows="Raw samples from the training split.",
                    setup="60 000 training images, 10 000 test.",
                    data="The MNIST training split as loaded.",
                    why="The project's contribution is the distributed learning setting, not "
                    "the vision task. A well-understood dataset keeps attention there.",
                ),
                dict(
                    file="04_downsampling.png",
                    title="Why 14×14",
                    claim="A quarter of the pixels, almost none of the accuracy.",
                    shows="The same digits at 28×28 and downsampled to 14×14.",
                    setup="Area-average downsampling, applied before normalisation.",
                    data="Training split.",
                    why="Sets p = 2908 parameters. Phase 5 needs a dense p×p covariance — "
                    "8.5 M entries per agent, which is affordable; at full resolution it "
                    "would not be. This choice is what makes the filter feasible.",
                ),
                dict(
                    file="05_rotation.png",
                    title="How the distribution drifts",
                    claim="Rotation is the drift mechanism: continuous, controllable, and it "
                    "preserves the label.",
                    shows="One digit at a range of rotations.",
                    setup="Rotation applied to the image before normalisation.",
                    data="Training split.",
                    why="Gives a non-stationarity with a single scalar state, so 'how far has "
                    "the world moved' has an exact answer at every step — which is what the "
                    "phase-5 state model will track.",
                ),
                dict(
                    file="07_drift_schedules.png",
                    title="Drift schedules",
                    claim="Four regimes: stationary, linear, piecewise, sinusoidal.",
                    shows="Rotation angle against step for each schedule.",
                    setup="Total rotation is capped at 45°, and the per-step rate α = 45/T is "
                    "**derived** rather than chosen.",
                    data="Constructed.",
                    why="Deriving α from T means the task stays equally hard regardless of "
                    "horizon — but it also means changing T changes the data at step t, which "
                    "is why runs at different horizons are not comparable.",
                ),
                dict(
                    file="08_partition_skew.png",
                    title="Non-IID data partitions",
                    claim="Dirichlet β controls how differently agents see the world.",
                    shows="Class composition per agent at several β.",
                    setup="Shard **sizes** are held equal; only the label composition varies.",
                    data="Training-split labels, partitioned at setup.",
                    why="Holding sizes fixed is what stops skew being confounded with shard "
                    "starvation. At β = 0.1 an agent sees three or four digits — the regime "
                    "where cooperation turns out to matter most.",
                ),
                dict(
                    file="10_received_digits.png",
                    title="What an agent actually receives",
                    claim="A few samples per step, from its own disjoint shard.",
                    shows="The samples delivered to each agent over the first steps of a run.",
                    setup="n = 4 samples per agent per step, label availability π = 1.",
                    data="One run of the real stream.",
                    why="Makes the online setting concrete: no agent ever holds the dataset, "
                    "and no sample is seen twice.",
                ),
                dict(
                    file="11_reference.png",
                    title="The offline reference e*",
                    claim="The error floor: what this architecture achieves given the full "
                    "shard and unlimited passes.",
                    shows="Test error against rotation for the offline classifier, per "
                    "initialisation strategy.",
                    setup="Same 196–14–10 MLP, 100 epochs, batch 128, epoch chosen on a 5 000 "
                    "image validation split, retrained per rotation.",
                    data="Full training split, offline.",
                    why="Same architecture, so the gap to it is purely 'offline versus online' "
                    "rather than model capacity. e* = 0.047 at 0°.",
                ),
            ],
        ),
        (
            "Results",
            "What the benchmark measured",
            [
                dict(
                    file="12_f1_error_vs_time.png",
                    title="F1 — the headline",
                    claim="Diffusion is statistically indistinguishable from pooled data.",
                    shows="Prequential error against t, stationary and rotating, with e* dashed.",
                    setup="Ring, N = 10, n = 4, T = 1500, every method at its own tuned rate.",
                    data="X1 and X2, five seeds, band ±1 s.d.",
                    why="The gap from ATC to centralized is 0.0014 against a seed s.d. of "
                    "0.0035 — inside the noise. This is the result that closes off stationary "
                    "accuracy as an avenue for phase 5.",
                ),
                dict(
                    file="13_f2_error_vs_communication.png",
                    title="F2 — error against bandwidth",
                    claim="The curves cross: cheapest is not the same as best.",
                    shows="F1 replotted against cumulative scalars transmitted, log x.",
                    setup="Centralized and local-only are horizontal references — they carry no "
                    "cost on this axis, for opposite reasons.",
                    data="X1 and X2, the communication ledger recorded per step.",
                    why="At equal *time* momentum ATC wins by 0.013; at equal *bandwidth* the "
                    "payload-matched variant wins until quite late. This is the axis on which "
                    "Diff-EKF's claim is actually stated.",
                ),
                dict(
                    file="16_f3_price_of_connectivity.png",
                    title="F3 — the price of connectivity",
                    claim="The spectral gap predicts, but the mean self-weight predicts better.",
                    shows="Gap to centralized against two graph quantities, in a transient and "
                    "a settled window.",
                    setup="Seven topologies, each at its own tuned learning rate. The gap is "
                    "paired within seed, which cuts the noise from 0.0035 to 0.0012.",
                    data="X3, five seeds per topology.",
                    why="A star has a *higher* spectral gap than a grid yet five times the "
                    "penalty. Mean self-weight ranks all seven at ρ = +0.964 (p = 0.0028) "
                    "because 1 − ā is literally mixing-per-round. Caveat: chosen post hoc, and "
                    "it failed its one out-of-sample test.",
                ),
                dict(
                    file="17_f4_per_agent_spread.png",
                    title="F4 — per-agent spread",
                    claim="How tightly the combine step holds the network together.",
                    shows="Mean over agents with a min–max band, per method.",
                    setup="Held-out error per agent, then the spread across agents.",
                    data="X1 and X2, five seeds.",
                    why="Centralized has no band at all — every agent identical by "
                    "construction, so a visible band there would mean a bug. Local-only's band "
                    "is wide. This is disagreement in *performance*, not in parameters.",
                ),
                dict(
                    file="14_f5_disagreement.png",
                    title="F5 — consensus and fidelity",
                    claim="The combine step turns a random walk into an equilibrium.",
                    shows="E_agree (how far agents are from each other), E_cent (how far their "
                    "average is from the centralized solution), and E_cent normalised by "
                    "‖θ̄‖² — all log y.",
                    setup="Computed on the parameter vectors, not on predictions. The third "
                    "row exists because the raw E_cent is easy to misread.",
                    data="X1 and X2, five seeds.",
                    why="ATC's E_agree plateaus at ~0.009 where the combine's pull toward "
                    "consensus balances the gradients pushing apart; local_only has no such "
                    "pull, so its disagreement grows linearly in t (r = 0.993) — a random "
                    "walk. E_cent rises for everyone, but about half of that is the weights "
                    "growing (‖θ̄‖² nearly doubles), and the rest costs almost nothing in "
                    "error (0.0749 vs 0.0762 at t = 1499): the models stay functionally "
                    "identical while their parameters drift along a flat direction.",
                ),
                dict(
                    file="18_f6a_sparsity_tuned.png",
                    title="F6a — the sparsity plane",
                    claim="Cooperation is worth eight times more when data is scarce.",
                    shows="ATC error, cooperation gap, and pooling gap over (n, π_lab).",
                    setup="288 tuning runs so every cell uses each method's own best rate. "
                    "Without that the comparison measures step size, not method.",
                    data="X4, T = 750, three seeds per cell.",
                    why="The cooperation gap runs 0.047 to 0.370 across the plane. The pooling "
                    "gap is *negative* only in the sparse corner and shrinks monotonically — "
                    "the signature of implicit iterate averaging, not of mis-tuning.",
                ),
                dict(
                    file="19_f6b_cost_of_not_retuning.png",
                    title="F6b — the cost of a wrong learning rate",
                    claim="Diffusion needs less re-tuning, because it re-tunes itself.",
                    shows="error(headline rate) − error(best rate for that cell), per method. "
                    "Worst penalty: centralized 0.217, local-only 0.144, ATC 0.027.",
                    setup="Only computable because X4 was run at both tunings — the fixed one "
                    "and the per-cell one.",
                    data="X4, both variants, three seeds.",
                    why="Not because ATC's error-vs-lr curve is flatter (it is not — 0.035 "
                    "against centralized's 0.037, and local_only is flattest of all yet "
                    "second-worst). It is that ATC's optimum barely moves: with ~2.5 of 10 "
                    "agents active, idle agents contribute unchanged θ to the combine, so its "
                    "effective step is η·n_active/N ≈ η/4 automatically — exactly the "
                    "reduction a smaller batch needs. Centralized applies the full η whatever "
                    "the batch, so its optimum has to move by 4× and the headline rate costs "
                    "it 0.18.",
                ),
                dict(
                    file="20_f7_adaptation_transient.png",
                    title="F7 — recovery from an abrupt shift",
                    claim="Every method loses the same amount, and recovers at the same rate.",
                    shows="Held-out error around a 15° jump at t = 500.",
                    setup="The change point is read from the recorded drift state, not "
                    "hardcoded. Deliberately unsmoothed — the spike lives in a single "
                    "evaluation point and any rolling window erases it.",
                    data="X5, five seeds.",
                    why="All three lose about 0.074 at once and recover within ~150 steps with "
                    "the ordering unchanged. The shift does not advantage any method — which "
                    "makes it a clean baseline for the filter to improve on.",
                ),
                dict(
                    file="15_f8_atc_vs_cta.png",
                    title="F8 — ATC vs CTA",
                    claim="Indistinguishable at tuned settings; ATC's advantage is robustness.",
                    shows="Prequential error and disagreement for the two orderings, at "
                    "identical communication cost.  **Prequential = test-then-train**: each "
                    "agent predicts on its batch **before** learning from it. Chosen because "
                    "it is honest online performance needing no held-out split, and because "
                    "scoring **after** the update leaks the label — by an amount growing with "
                    "the learning rate, so it would look like a result.",
                    setup="Both select the same optimum (momentum, lr 0.01), so matched "
                    "settings are also each one's best — checked, not assumed.",
                    data="X1b, five seeds, plus a 50-cell tuning grid.",
                    why="0.0768 vs 0.0774 against a s.d. of 0.0034. ATC wins 47 of 50 grid "
                    "cells so the sign is real, but it separates only where the step is too "
                    "large. ATC stays primary because Diff-EKF is ATC.",
                ),
                dict(
                    file="22_f10_forgetting.png",
                    title="F10 — forgetting",
                    claim="Nobody forgets. The lone agent lags.",
                    shows="Current vs backward error, and their paired gap. `backward` scores "
                    "the model at a rotation it visited earlier and has since left.",
                    setup="Sinusoidal rotation, amplitude 30°, period 500 — the only schedule "
                    "that revisits states, so the only one where forgetting is well posed. The "
                    "probe is defined for 97% of steps here, 67% under linear, 0% stationary.",
                    data="X7, five seeds. The `backward` evalset must be enabled explicitly.",
                    why="⚠ Read the dashed cycle means, not the peaks: the instantaneous gap "
                    "swings ±0.05 with the drift phase, and averaging over a fifth of a period "
                    "flips its sign (+0.016 vs −0.0035). No cooperative method forgets "
                    "measurably (1.1–1.6 σ). local_only is negative at 10 σ — better on a state "
                    "it left than the current one, which is lag, not retention.",
                ),
                dict(
                    file="21_f9_non_iid.png",
                    title="F9 — non-IID",
                    claim="The strongest result in the benchmark.",
                    shows="Error and cooperation gap against Dirichlet β, log x.",
                    setup="Shard sizes equal; only label composition varies.",
                    data="X6, T = 1500, three seeds.",
                    why="Cooperation goes from worth 0.061 to worth 0.518 — 8.5×. At β = 0.1 a "
                    "lone agent lands near chance (0.62) while the same agent in the network "
                    "reaches 0.106. Centralized is flat across β, which is a free check that "
                    "the skew is in the partition and not leaking into the data path.",
                ),
            ],
        ),
    ]


# --------------------------------------------------------------------------- #
# layout
# --------------------------------------------------------------------------- #


def blank(presentation: Presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = SURFACE
    return slide


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def write(frame, text, size, colour=INK, bold=False, space_after=4, first=False):
    """One paragraph. `**spans**` are set bold, so a note can emphasise a term.

    Without this the markers render as literal asterisks -- which is what
    happened the first time a note tried to emphasise something.
    """
    paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
    paragraph.space_after = Pt(space_after)
    for index, chunk in enumerate(text.split("**")):
        if not chunk:
            continue
        run = paragraph.add_run()
        run.text = chunk
        run.font.size = Pt(size)
        run.font.bold = bold or index % 2 == 1
        run.font.color.rgb = colour
        run.font.name = "Calibri"
    return paragraph


def fitted(path: Path, box_w: int, box_h: int) -> tuple[int, int]:
    """Largest (w, h) inside the box that preserves the image's aspect ratio."""
    with Image.open(path) as image:
        ratio = image.size[0] / image.size[1]
    if box_w / box_h > ratio:
        return int(box_h * ratio), box_h
    return box_w, int(box_w / ratio)


def rule(slide, top, width):
    line = slide.shapes.add_shape(1, MARGIN, top, width, Emu(9525))  # 1 = rectangle
    line.fill.solid()
    line.fill.fore_color.rgb = RULE
    line.line.fill.background()
    line.shadow.inherit = False


def figure_slide(presentation, spec: Slide, index: int, total: int) -> None:
    slide = blank(presentation)
    path = FIGURES / spec["file"]
    if not path.is_file():
        raise FileNotFoundError(path)

    with Image.open(path) as image:
        aspect = image.size[0] / image.size[1]

    frame = textbox(slide, MARGIN, Inches(0.32), WIDE - 2 * MARGIN, Inches(0.9))
    write(frame, spec["title"], 26, INK, bold=True, first=True, space_after=2)
    write(frame, spec["claim"], 14, ACCENT, space_after=0)

    top = Inches(1.5)
    rule(slide, Inches(1.38), WIDE - 2 * MARGIN)

    notes = [
        ("Shows", spec["shows"]),
        ("Setup", spec["setup"]),
        ("Data", spec["data"]),
        ("Why it matters", spec["why"]),
    ]

    if aspect >= 1.95:
        # Wide: image across the top, notes in two columns underneath.
        #
        # The image height is whatever the notes leave, rather than a fixed cap.
        # A cap sized for the wordiest slide wastes an inch of height on every
        # other one, and in a meeting the figure is the thing people need to be
        # able to read from the back of the room.
        column_in = (WIDE - 2 * MARGIN - Inches(0.4)) / 2 / 914400

        # Split the four notes across the two columns to *balance* them, rather
        # than 2 + 2. One long note (a slide that explains a term as well as a
        # result) otherwise makes its column tall while the other sits half
        # empty, and since the image gets whatever height the notes leave, the
        # figure pays for the imbalance.
        best_split, best_cost = 2, None
        for cut in (1, 2, 3):
            left = sum(len(text) for _, text in notes[:cut]) + 60 * cut
            right = sum(len(text) for _, text in notes[cut:]) + 60 * (len(notes) - cut)
            cost = max(left, right)
            if best_cost is None or cost < best_cost:
                best_split, best_cost = cut, cost

        note_lines = best_cost / (column_in * 9.0) + 2
        notes_in = min(2.6, max(1.4, note_lines * 0.19))
        available = 7.5 - 1.5 - notes_in - 0.22 - 0.30
        box_w, box_h = int(WIDE - 2 * MARGIN), Inches(max(2.8, available))
        width, height = fitted(path, box_w, box_h)
        slide.shapes.add_picture(
            str(path), int((WIDE - width) / 2), top, width=width, height=height
        )
        notes_top = top + height + Inches(0.22)
        column = int((WIDE - 2 * MARGIN - Inches(0.4)) / 2)
        for offset, chunk in (
            (0, notes[:best_split]),
            (column + Inches(0.4), notes[best_split:]),
        ):
            frame = textbox(
                slide, MARGIN + offset, notes_top, column, TALL - notes_top - Inches(0.3)
            )
            for i, (label, text) in enumerate(chunk):
                write(frame, label.upper(), 9.5, MUTED, bold=True, first=(i == 0), space_after=1)
                write(frame, text, 11.5, INK_SECONDARY, space_after=7)
    else:
        # Tall or square: image left, notes in a right-hand column.
        box_w, box_h = int(WIDE * 0.56), int(TALL - top - Inches(0.35))
        width, height = fitted(path, box_w, box_h)
        slide.shapes.add_picture(str(path), MARGIN, top, width=width, height=height)
        left = MARGIN + box_w + Inches(0.3)
        frame = textbox(slide, left, top, WIDE - left - MARGIN, TALL - top - Inches(0.35))
        for i, (label, text) in enumerate(notes):
            write(frame, label.upper(), 9.5, MUTED, bold=True, first=(i == 0), space_after=1)
            write(frame, text, 11.5, INK_SECONDARY, space_after=9)

    frame = textbox(slide, WIDE - Inches(1.3), TALL - Inches(0.42), Inches(0.9), Inches(0.3))
    paragraph = write(frame, f"{index} / {total}", 9, MUTED, first=True, space_after=0)
    paragraph.alignment = PP_ALIGN.RIGHT


def title_slide(presentation) -> None:
    slide = blank(presentation)
    frame = textbox(slide, MARGIN, Inches(2.3), WIDE - 2 * MARGIN, Inches(3))
    write(
        frame,
        "A Benchmark for Distributed Online Learning over a Graph",
        36,
        INK,
        bold=True,
        first=True,
        space_after=10,
    )
    write(frame, "Phases 0–4 — groundwork for the diffusion EKF", 18, ACCENT, space_after=22)
    write(frame, "Tomer Alter · Ben-Gurion University", 14, INK_SECONDARY, space_after=4)
    write(
        frame,
        "10 agents · one shared classifier · no fusion centre · "
        "196–14–10 MLP, p = 2908 · MNIST at 14×14",
        12,
        MUTED,
        space_after=0,
    )


def section_slide(presentation, title: str, subtitle: str) -> None:
    slide = blank(presentation)
    frame = textbox(slide, MARGIN, Inches(3.0), WIDE - 2 * MARGIN, Inches(1.6))
    write(frame, title, 32, INK, bold=True, first=True, space_after=8)
    write(frame, subtitle, 15, INK_SECONDARY, space_after=0)


def closing_slide(presentation) -> None:
    slide = blank(presentation)
    frame = textbox(slide, MARGIN, Inches(0.6), WIDE - 2 * MARGIN, Inches(1.0))
    write(frame, "Where phase 5 can make its case", 28, INK, bold=True, first=True, space_after=3)
    write(
        frame,
        "The benchmark closed one avenue and left five open, plus two worth discussing.",
        14,
        ACCENT,
        space_after=0,
    )
    rule(slide, Inches(1.65), WIDE - 2 * MARGIN)

    rows = [
        (
            "CLOSED",
            "Stationary accuracy",
            "ATC is inside seed noise of pooled data (0.0014 vs s.d. 0.0035); the worst "
            "settled topology costs 0.008. There is no gap left to close.",
        ),
        (
            "open",
            "Tracking under drift",
            "Gaps are wider and still opening under rotation and after abrupt shifts.",
        ),
        (
            "open",
            "Communication efficiency",
            "Must beat ATC (payload-matched) at 0.0902 — not 0.0768 — at p per link.",
        ),
        ("open", "Calibrated uncertainty", "No SGD baseline provides it at all."),
        (
            "open",
            "Information-weighted combine",
            "Fixed weights ignore how much data a neighbour held — worst at π = 0.25 and β = 0.1, "
            "exactly where the gaps are largest.",
        ),
        (
            "open",
            "Curvature-derived step",
            "No learning rate to get wrong (F6b: 0.217 vs 0.027) — though a prior scale and a "
            "forgetting factor replace it.",
        ),
        (
            "exploratory",
            "Drift detection",
            "The innovation covariance gives a test statistic — but no baseline attempts it, so "
            "there is nothing to compare against.",
        ),
        (
            "measured",
            "Controlled forgetting",
            "X7 now runs it: no cooperative method forgets measurably (1.1–1.6 σ), and "
            "local_only lags at 10 σ. So a filter cannot win by fixing forgetting — but λ "
            "trades tracking against retention, and X7 is where that trade becomes visible.",
        ),
    ]
    top = Inches(1.95)
    frame = textbox(slide, MARGIN, top, WIDE - 2 * MARGIN, TALL - top - Inches(0.4))
    for i, (status, avenue, text) in enumerate(rows):
        colour = INK if status == "CLOSED" else (MUTED if status == "exploratory" else ACCENT)
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(7)
        for content, size, bold, rgb in (
            (f"{status:<12}", 10.5, True, colour),
            (f"{avenue}  ", 12.5, True, INK),
            (text, 11.5, False, INK_SECONDARY),
        ):
            run = paragraph.add_run()
            run.text = content
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb
            run.font.name = "Calibri"


def build() -> Path:
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = WIDE, TALL

    title_slide(presentation)
    groups = sections()
    total = sum(len(slides) for _, _, slides in groups)

    index = 0
    for title, subtitle, slides in groups:
        section_slide(presentation, title, subtitle)
        for spec in slides:
            index += 1
            figure_slide(presentation, spec, index, total)
    closing_slide(presentation)

    path = FIGURES / "Phase_0-4_Review.pptx"
    presentation.save(path)
    return path


if __name__ == "__main__":
    written = build()
    print(f"wrote {written}")
    print(f"  {len(Presentation(written).slides._sldIdLst)} slides")
